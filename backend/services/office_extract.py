from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
import shutil
import subprocess
from typing import Iterable
from uuid import uuid4

from docx import Document
from openpyxl import load_workbook
from pdfminer.high_level import extract_text as pdf_extract_text
from pptx import Presentation

from backend.i18n import get_list

ALLOWED_OFFICE_EXTENSIONS = {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".pdf"}
TITLE_STYLE_ALIASES = {alias.lower() for alias in get_list("office.title_styles")}


@dataclass
class ExtractionResult:
    markdown: str
    warnings: list[str]
    source_type: str


def extract_office_to_markdown(input_path: Path) -> ExtractionResult:
    suffix = input_path.suffix.lower()
    if suffix == ".docx":
        return _docx_to_markdown(input_path)
    if suffix == ".xlsx":
        return _xlsx_to_markdown(input_path)
    if suffix == ".pptx":
        return _pptx_to_markdown(input_path)
    if suffix == ".pdf":
        return _pdf_to_markdown(input_path)
    if suffix in {".doc", ".xls", ".ppt"}:
        return _legacy_to_markdown(input_path, suffix)
    raise ValueError(f"Unsupported office file type: {suffix}")


def _clean_text(value: object) -> str:
    text = str(value or "").strip()
    text = re.sub(r"\s+", " ", text)
    return text


def _docx_to_markdown(path: Path) -> ExtractionResult:
    doc = Document(str(path))
    lines: list[str] = []
    warnings: list[str] = []

    for para in doc.paragraphs:
        text = _clean_text(para.text)
        if not text:
            continue
        style_name = getattr(getattr(para, "style", None), "name", "") or ""
        if style_name.startswith("Heading"):
            level_text = style_name.replace("Heading", "").strip()
            level = int(level_text) if level_text.isdigit() else 1
            level = max(1, min(level, 6))
            lines.append(f"{'#' * level} {text}")
        elif style_name.lower() in TITLE_STYLE_ALIASES:
            lines.append(f"# {text}")
        else:
            lines.append(text)
        lines.append("")

    if doc.tables:
        warnings.append("Tables appended after paragraphs; original order in document may differ.")
        for idx, table in enumerate(doc.tables, start=1):
            lines.append(f"## Table {idx}")
            rows = list(table.rows)
            if not rows:
                continue
            first = [_clean_text(cell.text) for cell in rows[0].cells]
            header = first if any(first) else [f"Column {i+1}" for i in range(len(first))]
            lines.append(_md_table_row(header))
            lines.append(_md_table_sep(len(header)))
            start_row = 1 if any(first) else 0
            for row in rows[start_row:]:
                cells = [_clean_text(cell.text) for cell in row.cells]
                lines.append(_md_table_row(cells))
            lines.append("")

    markdown = "\n".join(lines).strip() + "\n"
    return ExtractionResult(markdown=markdown, warnings=warnings, source_type="docx")


def _xlsx_to_markdown(path: Path) -> ExtractionResult:
    wb = load_workbook(filename=str(path), data_only=True, read_only=True)
    lines: list[str] = []
    warnings: list[str] = []

    max_rows = 200
    max_cols = 50

    for sheet in wb.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        rows = []
        for i, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if i > max_rows:
                warnings.append(f"Sheet '{sheet.title}' truncated to {max_rows} rows.")
                break
            rows.append(list(row[:max_cols]))

        if not rows:
            lines.append("(empty sheet)")
            lines.append("")
            continue

        # Determine column count by last non-empty cell in any row.
        col_count = 0
        for row in rows:
            for idx in range(len(row) - 1, -1, -1):
                if row[idx] not in (None, ""):
                    col_count = max(col_count, idx + 1)
                    break
        col_count = max(col_count, 1)

        def row_values(raw: Iterable[object]) -> list[str]:
            values = list(raw)[:col_count]
            if len(values) < col_count:
                values.extend([""] * (col_count - len(values)))
            return [_clean_text(v) for v in values]

        header = row_values(rows[0])
        if not any(header):
            header = [f"Column {i+1}" for i in range(col_count)]
            start_idx = 0
        else:
            start_idx = 1

        lines.append(_md_table_row(header))
        lines.append(_md_table_sep(len(header)))
        for row in rows[start_idx:]:
            lines.append(_md_table_row(row_values(row)))
        lines.append("")

    markdown = "\n".join(lines).strip() + "\n"
    return ExtractionResult(markdown=markdown, warnings=warnings, source_type="xlsx")


def _pptx_to_markdown(path: Path) -> ExtractionResult:
    prs = Presentation(str(path))
    lines: list[str] = []
    warnings: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title_text = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = _clean_text(slide.shapes.title.text_frame.text)
        heading = f"# Slide {idx}"
        if title_text:
            heading = f"{heading}: {title_text}"
        lines.append(heading)

        bullets: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = _clean_text(paragraph.text)
                if not text or text == title_text:
                    continue
                level = getattr(paragraph, "level", 0) or 0
                indent = "  " * min(level, 5)
                bullets.append(f"{indent}- {text}")

        if bullets:
            lines.extend(bullets)
        else:
            warnings.append(f"Slide {idx} has no extractable text.")
        lines.append("")

    markdown = "\n".join(lines).strip() + "\n"
    return ExtractionResult(markdown=markdown, warnings=warnings, source_type="pptx")


def _pdf_to_markdown(path: Path) -> ExtractionResult:
    warnings: list[str] = []
    text = pdf_extract_text(str(path)) or ""
    text = text.replace("\r", "")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        warnings.append("No extractable text found in PDF.")
    max_chars = 200_000
    truncated = len(text) > max_chars
    if truncated:
        text = text[:max_chars]
        warnings.append(f"PDF content truncated to {max_chars} characters.")
    markdown = (text + "\n") if text else ""
    return ExtractionResult(markdown=markdown, warnings=warnings, source_type="pdf")


def _legacy_to_markdown(path: Path, suffix: str) -> ExtractionResult:
    if suffix == ".doc":
        target_ext = ".docx"
    elif suffix == ".xls":
        target_ext = ".xlsx"
    elif suffix == ".ppt":
        target_ext = ".pptx"
    else:
        raise ValueError(f"Unsupported legacy format: {suffix}")

    converted = _convert_with_soffice(path, target_ext)
    result = extract_office_to_markdown(converted)
    warnings = ["Converted with LibreOffice for legacy format support."] + result.warnings
    return ExtractionResult(markdown=result.markdown, warnings=warnings, source_type=suffix.lstrip("."))


def _convert_with_soffice(path: Path, target_ext: str) -> Path:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) is required to convert legacy Office formats.")

    out_dir = path.parent
    tmp_dir = out_dir / f".convert_{uuid4().hex[:8]}"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    try:
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                target_ext.lstrip("."),
                "--outdir",
                str(tmp_dir),
                str(path),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
        candidate = tmp_dir / f"{path.stem}{target_ext}"
        if not candidate.exists():
            outputs = list(tmp_dir.glob(f"*{target_ext}"))
            if outputs:
                candidate = outputs[0]
        if not candidate.exists():
            raise RuntimeError("LibreOffice conversion produced no output file.")

        final_path = out_dir / f"{path.stem}_converted_{uuid4().hex[:6]}{target_ext}"
        candidate.replace(final_path)
        return final_path
    finally:
        try:
            for item in tmp_dir.iterdir():
                item.unlink(missing_ok=True)
            tmp_dir.rmdir()
        except Exception:
            pass


def _md_table_row(values: Iterable[str]) -> str:
    escaped = [str(v).replace("|", "\\|") for v in values]
    return "| " + " | ".join(escaped) + " |"


def _md_table_sep(count: int) -> str:
    return "| " + " | ".join(["---"] * max(1, count)) + " |"
